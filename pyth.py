from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

# Create presentation
prs = Presentation()

# UN / SDG Colors
SDG_BLUE = RGBColor(0, 122, 193)
SDG_GREEN = RGBColor(76, 175, 80)
SDG_ORANGE = RGBColor(255, 152, 0)
SDG_PURPLE = RGBColor(156, 39, 176)
SDG_GRAY = RGBColor(220, 220, 220)
WHITE = RGBColor(255, 255, 255)

def title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.3), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = SDG_BLUE

    box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf = box.text_frame
    tf.text = "GLOBAL IMPACT AWARDS 2025"
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Plan-et Planners | UN SDG–Aligned Global Award Ceremony\nDecember 20, 2025"
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

def bullet_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = SDG_BLUE
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            tf.text = b
        else:
            tf.add_paragraph().text = b

# ---- Slides 1–10 ----
title_slide()

bullet_slide("Event Overview", [
    "Global Impact Awards recognizes organizations advancing the UN SDGs.",
    "Celebrates innovation, sustainability, ethical leadership, and global responsibility.",
    "25 nominated organizations with over 250 international delegates.",
    "4-hour premium award ceremony on December 20, 2025."
])

bullet_slide("Vision & Purpose", [
    "Promote sustainable global business practices.",
    "Encourage cross-border SDG collaboration.",
    "Position sustainability as a competitive advantage.",
    "Inspire innovation through recognition."
])

bullet_slide("Target Audience", [
    "CEOs and business leaders.",
    "NGOs and development partners.",
    "Government and policymakers.",
    "Global media and investors."
])

bullet_slide("Award Categories", [
    "Best SDG Innovation Award.",
    "Sustainability Leadership Award.",
    "Global Social Impact Award.",
    "Climate Action Excellence Award.",
    "Ethical Business Practices Award."
])

bullet_slide("Client Consultation", [
    "Define event objectives and KPIs.",
    "Align with SDG priorities.",
    "Confirm categories and judging criteria.",
    "Finalize event scale and tone."
])

bullet_slide("Marketing & Promotion", [
    "Digital advertising campaigns.",
    "Influencer partnerships.",
    "Email marketing.",
    "Press releases and SDG storytelling."
])

bullet_slide("Nomination & Voting", [
    "Official nomination notifications.",
    "Online voting system.",
    "Category-based scoring.",
    "Independent review for transparency."
])

bullet_slide("Partnerships & Sponsorships", [
    "UN-aligned sponsors.",
    "Media partnerships.",
    "Ethical corporate brands.",
    "Shared sustainability values."
])

bullet_slide("Venue & Logistics", [
    "300-seat premium hall.",
    "Stage, lighting, sound, live streaming.",
    "Catering, cooling, security.",
    "Transportation and coordination."
])

# ---- Slide 11: Floor Plan ----
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Venue Floor Plan – Award Ceremony Layout"
desc = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(3))
dtf = desc.text_frame
dtf.text = (
    "The venue layout ensures smooth guest movement, visibility, and security.\n\n"
    "• Registration desk at entrance\n"
    "• VIP seating closest to stage\n"
    "• Central audience seating (250 guests)\n"
    "• Dedicated media zone\n"
    "• Technical control area"
)

stage = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(1.5), Inches(5), Inches(0.8))
stage.fill.solid(); stage.fill.fore_color.rgb = SDG_PURPLE
stage.text_frame.text = "MAIN STAGE"

vip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(2.5), Inches(5), Inches(1))
vip.fill.solid(); vip.fill.fore_color.rgb = SDG_ORANGE
vip.text_frame.text = "VIP SEATING"

seat = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(3.7), Inches(5), Inches(1.5))
seat.fill.solid(); seat.fill.fore_color.rgb = SDG_GRAY
seat.text_frame.text = "GENERAL AUDIENCE SEATING"

reg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(5.5), Inches(2.4), Inches(0.8))
reg.fill.solid(); reg.fill.fore_color.rgb = SDG_GREEN
reg.text_frame.text = "REGISTRATION"

media = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.6), Inches(5.5), Inches(2.4), Inches(0.8))
media.fill.solid(); media.fill.fore_color.rgb = SDG_BLUE
media.text_frame.text = "MEDIA ZONE"

# ---- Slide 12: Timeline ----
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Event Timeline"

timeline = [
    ("Consultation", "Jan 2025"),
    ("Marketing", "Feb–Mar 2025"),
    ("Nominations", "Apr–May 2025"),
    ("Final Prep", "Nov 2025"),
    ("Event Day", "Dec 20, 2025")
]

for i, (p, d) in enumerate(timeline):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON, Inches(0.5 + i*2.5), Inches(3), Inches(2.3), Inches(1)
    )
    sh.fill.solid(); sh.fill.fore_color.rgb = SDG_BLUE
    tf = sh.text_frame
    tf.text = p
    q = tf.add_paragraph()
    q.text = d
    tf.paragraphs[0].font.color.rgb = WHITE
    q.font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    q.alignment = PP_ALIGN.CENTER

# ---- Slide 13: Budget ----
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Estimated Budget Allocation (₦)"

chart_data = ChartData()
chart_data.categories = [
    "Venue", "Catering", "Technology", "Marketing",
    "Awards", "Logistics", "Security", "Staffing"
]
chart_data.add_series(
    "Budget",
    (800000, 1200000, 900000, 600000, 400000, 300000, 500000, 500000)
)

slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(1), Inches(2), Inches(11), Inches(4),
    chart_data
)

# ---- Slides 14–15 ----
bullet_slide("Risk Management", [
    "Backup power and cooling systems.",
    "Professional security and medical team.",
    "Contingency planning.",
    "Real-time monitoring."
])

bullet_slide("Expected Outcomes & Legacy", [
    "Recognition of impactful organizations.",
    "Stronger SDG partnerships.",
    "Enhanced Plan-et Planners reputation.",
    "Framework for future global events."
])

# Save file
prs.save("Global_Impact_Awards_2025_Final_15_Slides_With_Floor_Plan.pptx")